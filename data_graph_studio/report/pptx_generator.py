"""
PowerPoint (PPTX) Report Generator
PowerPoint 프레젠테이션 레포트 생성기

Uses python-pptx for presentation generation.
Creates presentation-ready slides with charts and data.
"""

import io
import base64
import logging

from data_graph_studio.core.report import (
    ReportGenerator,
    ReportData,
    ReportOptions,
)

logger = logging.getLogger(__name__)

# Check if python-pptx is available
PPTX_AVAILABLE = False
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu  # noqa: F401
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: F401
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR  # noqa: F401
    PPTX_AVAILABLE = True
except ImportError:
    pass


class PPTXReportGenerator(ReportGenerator):
    """PowerPoint 프레젠테이션 생성기"""

    # Slide sizes (lazy loaded)
    _SLIDE_SIZES = None
    
    @classmethod
    def _get_slide_sizes(cls):
        if cls._SLIDE_SIZES is None and PPTX_AVAILABLE:
            cls._SLIDE_SIZES = {
                "16:9": (Inches(13.333), Inches(7.5)),
                "4:3": (Inches(10), Inches(7.5)),
            }
        return cls._SLIDE_SIZES or {}

    def generate(
        self,
        report_data: ReportData,
        options: ReportOptions
    ) -> bytes:
        """PowerPoint 프레젠테이션 생성"""
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is not available. Install it with:\n"
                "  pip install python-pptx"
            )

        prs = Presentation()

        # 슬라이드 크기 설정
        slide_size = self.SLIDE_SIZES.get(options.slide_size, self.SLIDE_SIZES["16:9"])
        prs.slide_width = slide_size[0]
        prs.slide_height = slide_size[1]

        # 타이틀 슬라이드
        self._add_title_slide(prs, report_data, options)

        # Executive Summary 슬라이드
        if options.include_executive_summary:
            self._add_executive_summary_slide(prs, report_data, options)

        # Data Overview 슬라이드
        if options.include_data_overview:
            self._add_data_overview_slide(prs, report_data, options)

        # Statistics 슬라이드
        if options.include_statistics:
            self._add_statistics_slides(prs, report_data, options)

        # Visualizations 슬라이드
        if options.include_visualizations and report_data.charts:
            self._add_visualization_slides(prs, report_data, options)

        # Comparison 슬라이드
        if options.include_comparison and report_data.is_multi_dataset():
            self._add_comparison_slides(prs, report_data, options)

        # Data Tables 슬라이드
        if options.include_tables and report_data.tables:
            self._add_table_slides(prs, report_data, options)

        # Thank You 슬라이드
        self._add_closing_slide(prs, report_data, options)

        # 바이트로 저장
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        return buffer.getvalue()

    def _hex_to_rgb(self, hex_color: str) -> "RGBColor":
        """HEX를 RGBColor로 변환"""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def _add_title_slide(
        self,
        prs: "Presentation",
        report_data: ReportData,
        options: ReportOptions
    ):
        """타이틀 슬라이드 추가"""
        # 빈 슬라이드 레이아웃 사용
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 배경색 (그라디언트 효과를 위해 shape 사용)
        width = prs.slide_width
        height = prs.slide_height

        # 배경 shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self._hex_to_rgb(self.template.primary_color)
        bg_shape.line.fill.background()

        # 로고 (있는 경우)
        if report_data.metadata.logo_base64:
            try:
                logo_bytes = base64.b64decode(report_data.metadata.logo_base64)
                logo_stream = io.BytesIO(logo_bytes)

                logo_left = (width - Inches(2)) / 2
                slide.shapes.add_picture(
                    logo_stream,
                    logo_left, Inches(1),
                    width=Inches(2)
                )
            except Exception as e:
                logger.warning("pptx_generator.add_logo_failed", extra={"error": e})

        # 제목
        title_left = Inches(0.5)
        title_top = Inches(2.5)
        title_width = width - Inches(1)
        title_height = Inches(1.5)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = report_data.metadata.title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 부제목
        if report_data.metadata.subtitle:
            subtitle_top = Inches(4)
            subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, title_width, Inches(0.7))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.paragraphs[0].text = report_data.metadata.subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 메타 정보
        meta_top = Inches(6)
        meta_box = slide.shapes.add_textbox(title_left, meta_top, title_width, Inches(0.8))
        meta_frame = meta_box.text_frame

        meta_text = report_data.metadata.created_at.strftime("%Y-%m-%d")
        if report_data.metadata.author:
            meta_text += f" | {report_data.metadata.author}"

        meta_frame.paragraphs[0].text = meta_text
        meta_frame.paragraphs[0].font.size = Pt(14)
        meta_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)
        meta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_section_slide(
        self,
        prs: "Presentation",
        title: str,
        options: ReportOptions
    ):
        """섹션 구분 슬라이드 추가"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        width = prs.slide_width
        height = prs.slide_height

        # 배경
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self._hex_to_rgb(self.template.secondary_color)
        bg_shape.line.fill.background()

        # 섹션 제목
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            (height - Inches(1)) / 2,
            width - Inches(1),
            Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return slide

    def _add_content_slide(
        self,
        prs: "Presentation",
        title: str,
        options: ReportOptions
    ) -> tuple:
        """콘텐츠 슬라이드 추가 (제목과 콘텐츠 영역 반환)"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        width = prs.slide_width
        height = prs.slide_height

        # 제목 영역
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            width - Inches(1), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.template.primary_color)

        # 하단 선
        line_top = Inches(1)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), line_top,
            width - Inches(1), Pt(3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.template.primary_color)
        line.line.fill.background()

        # 콘텐츠 영역 정보 반환
        content_area = {
            'left': Inches(0.5),
            'top': Inches(1.3),
            'width': width - Inches(1),
            'height': height - Inches(1.8)
        }

        return slide, content_area

    def _add_executive_summary_slide(
        self,
        prs: "Presentation",
        report_data: ReportData,
        options: ReportOptions
    ):
        """Executive Summary 슬라이드"""
        is_ko = options.language == 'ko'
        title = "요약" if is_ko else "Executive Summary"

        slide, content_area = self._add_content_slide(prs, title, options)

        # 메트릭 카드
        metrics = [
            ("Total Rows" if not is_ko else "총 행 수",
             self.format_number(report_data.get_total_rows(), 0)),
            ("Datasets" if not is_ko else "데이터셋",
             str(len(report_data.datasets))),
        ]

        if report_data.datasets:
            metrics.append((
                "Columns" if not is_ko else "컬럼 수",
                str(report_data.datasets[0].column_count)
            ))

        if report_data.charts:
            metrics.append((
                "Charts" if not is_ko else "차트",
                str(len(report_data.charts))
            ))

        # 메트릭 카드 그리기
        card_width = Inches(2.5)
        card_height = Inches(1.5)
        card_spacing = Inches(0.3)

        num_cards = len(metrics)
        total_width = (num_cards * card_width) + ((num_cards - 1) * card_spacing)
        start_left = content_area['left'] + (content_area['width'] - total_width) / 2

        for i, (label, value) in enumerate(metrics):
            card_left = start_left + i * (card_width + card_spacing)
            card_top = content_area['top'] + Inches(0.3)

            # 카드 배경
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                card_left, card_top,
                card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(245, 247, 250)
            card.line.color.rgb = RGBColor(220, 225, 230)

            # 라벨
            label_box = slide.shapes.add_textbox(
                card_left, card_top + Inches(0.2),
                card_width, Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.paragraphs[0].text = label
            label_frame.paragraphs[0].font.size = Pt(12)
            label_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)
            label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 값
            value_box = slide.shapes.add_textbox(
                card_left, card_top + Inches(0.6),
                card_width, Inches(0.7)
            )
            value_frame = value_box.text_frame
            value_frame.paragraphs[0].text = value
            value_frame.paragraphs[0].font.size = Pt(32)
            value_frame.paragraphs[0].font.bold = True
            value_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(self.template.primary_color)
            value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Key Findings (있는 경우)
        if report_data.key_findings:
            findings_top = content_area['top'] + Inches(2.2)
            findings_title = "핵심 발견 사항" if is_ko else "Key Findings"

            title_box = slide.shapes.add_textbox(
                content_area['left'], findings_top,
                content_area['width'], Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.paragraphs[0].text = findings_title
            title_frame.paragraphs[0].font.size = Pt(18)
            title_frame.paragraphs[0].font.bold = True

            # 발견 사항 목록
            list_top = findings_top + Inches(0.5)
            list_box = slide.shapes.add_textbox(
                content_area['left'] + Inches(0.3), list_top,
                content_area['width'] - Inches(0.6), Inches(3)
            )
            list_frame = list_box.text_frame
            list_frame.word_wrap = True

            for i, finding in enumerate(report_data.key_findings[:5]):
                if i == 0:
                    p = list_frame.paragraphs[0]
                else:
                    p = list_frame.add_paragraph()
                p.text = f"• {finding}"
                p.font.size = Pt(14)
                p.space_before = Pt(6)

    def _add_data_overview_slide(
        self,
        prs: "Presentation",
        report_data: ReportData,
        options: ReportOptions
    ):
        """Data Overview 슬라이드"""
        is_ko = options.language == 'ko'
        title = "데이터 개요" if is_ko else "Data Overview"

        slide, content_area = self._add_content_slide(prs, title, options)

        # 데이터셋 카드
        num_datasets = len(report_data.datasets)
        card_width = min(Inches(4), (content_area['width'] - Inches(0.5)) / min(num_datasets, 3))
        card_height = Inches(2.5)

        for i, ds in enumerate(report_data.datasets[:3]):  # 최대 3개
            col = i % 3
            card_left = content_area['left'] + col * (card_width + Inches(0.25))
            card_top = content_area['top'] + Inches(0.3)

            # 카드 배경
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                card_left, card_top,
                card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = self._hex_to_rgb(ds.color)
            card.line.width = Pt(3)

            # 데이터셋 이름
            name_box = slide.shapes.add_textbox(
                card_left + Inches(0.15), card_top + Inches(0.15),
                card_width - Inches(0.3), Inches(0.5)
            )
            name_frame = name_box.text_frame
            name_frame.paragraphs[0].text = ds.name
            name_frame.paragraphs[0].font.size = Pt(16)
            name_frame.paragraphs[0].font.bold = True
            name_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(ds.color)

            # 데이터셋 정보
            info_lines = [
                f"Rows: {self.format_number(ds.row_count, 0)}",
                f"Columns: {ds.column_count}",
                f"Memory: {self.format_bytes(ds.memory_bytes)}",
            ]

            info_box = slide.shapes.add_textbox(
                card_left + Inches(0.15), card_top + Inches(0.7),
                card_width - Inches(0.3), Inches(1.5)
            )
            info_frame = info_box.text_frame

            for j, line in enumerate(info_lines):
                if j == 0:
                    p = info_frame.paragraphs[0]
                else:
                    p = info_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(12)
                p.space_before = Pt(4)

    def _add_statistics_slides(
        self,
        prs: "Presentation",
        report_data: ReportData,
        options: ReportOptions
    ):
        """Statistics 슬라이드들"""
        is_ko = options.language == 'ko'

        for dataset_id, stats_list in report_data.statistics.items():
            if not stats_list:
                continue

            # 데이터셋 이름 찾기
            dataset_name = dataset_id
            for ds in report_data.datasets:
                if ds.id == dataset_id:
                    dataset_name = ds.name
                    break

            title = f"통계 요약: {dataset_name}" if is_ko else f"Statistics: {dataset_name}"
            slide, content_area = self._add_content_slide(prs, title, options)

            # 테이블 생성
            headers = ['Column', 'Mean', 'Median', 'Std', 'Min', 'Max']
            num_cols = len(headers)
            num_rows = min(len(stats_list), 10) + 1  # 헤더 + 최대 10행

            table_width = content_area['width'] - Inches(0.5)
            table_width / num_cols
            row_height = Inches(0.4)

            table = slide.shapes.add_table(
                num_rows, num_cols,
                content_area['left'] + Inches(0.25),
                content_area['top'] + Inches(0.3),
                table_width,
                row_height * num_rows
            ).table

            # 헤더 스타일
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._hex_to_rgb(self.template.primary_color)

                para = cell.text_frame.paragraphs[0]
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.size = Pt(11)
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER

            # 데이터 행
            for row_idx, stat in enumerate(stats_list[:10]):
                values = [
                    stat.column[:15],
                    self.format_number(stat.mean, 2),
                    self.format_number(stat.median, 2),
                    self.format_number(stat.std, 2),
                    self.format_number(stat.min, 2),
                    self.format_number(stat.max, 2),
                ]

                for col_idx, value in enumerate(values):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = value

                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(10)
                    para.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

    def _add_visualization_slides(
        self,
        prs: "Presentation",
        report_data: ReportData,
        options: ReportOptions
    ):
        """Visualization 슬라이드들"""
        for chart in report_data.charts:
            slide, content_area = self._add_content_slide(prs, chart.title, options)

            # 차트 이미지 삽입
            if chart.image_bytes or chart.image_base64:
                try:
                    if chart.image_bytes:
                        image_stream = io.BytesIO(chart.image_bytes)
                    else:
                        image_bytes = base64.b64decode(chart.image_base64)
                        image_stream = io.BytesIO(image_bytes)

                    # 이미지 크기 계산 (콘텐츠 영역에 맞춤)
                    img_width = content_area['width'] - Inches(1)
                    img_height = content_area['height'] - Inches(0.5)

                    # 비율 유지
                    aspect_ratio = chart.width / chart.height if chart.height > 0 else 1.33
                    if img_width / img_height > aspect_ratio:
                        img_width = img_height * aspect_ratio
                    else:
                        img_height = img_width / aspect_ratio

                    img_left = content_area['left'] + (content_area['width'] - img_width) / 2
                    img_top = content_area['top'] + Inches(0.25)

                    slide.shapes.add_picture(
                        image_stream,
                        img_left, img_top,
                        width=img_width
                    )

                except Exception as e:
                    logger.warning("pptx_generator.add_chart_image_failed", extra={"error": e})

                    # 플레이스홀더 텍스트
                    text_box = slide.shapes.add_textbox(
                        content_area['left'], content_area['top'] + Inches(2),
                        content_area['width'], Inches(1)
                    )
                    text_frame = text_box.text_frame
                    text_frame.paragraphs[0].text = "[Chart image not available]"
                    text_frame.paragraphs[0].font.italic = True
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 차트 설명 (있는 경우)
            if chart.description:
                desc_top = content_area['top'] + content_area['height'] - Inches(0.5)
                desc_box = slide.shapes.add_textbox(
                    content_area['left'], desc_top,
                    content_area['width'], Inches(0.4)
                )
                desc_frame = desc_box.text_frame
                desc_frame.paragraphs[0].text = chart.description
                desc_frame.paragraphs[0].font.size = Pt(11)
                desc_frame.paragraphs[0].font.italic = True
                desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_comparison_slides(
        self,
        prs: "Presentation",
        report_data: ReportData,
        options: ReportOptions
    ):
        """Comparison 슬라이드들"""
        is_ko = options.language == 'ko'

        # Statistical Test Results
        if report_data.comparisons:
            title = "통계 검정 결과" if is_ko else "Statistical Test Results"
            slide, content_area = self._add_content_slide(prs, title, options)

            current_top = content_area['top'] + Inches(0.2)

            for comp in report_data.comparisons[:3]:  # 최대 3개
                # 검정 제목
                test_title = f"{comp.test_type}: {comp.dataset_a_name} vs {comp.dataset_b_name}"

                title_box = slide.shapes.add_textbox(
                    content_area['left'], current_top,
                    content_area['width'], Inches(0.4)
                )
                title_frame = title_box.text_frame
                title_frame.paragraphs[0].text = test_title
                title_frame.paragraphs[0].font.size = Pt(14)
                title_frame.paragraphs[0].font.bold = True

                current_top += Inches(0.45)

                # 결과 정보
                sig_text = "Not Significant"
                if comp.p_value < 0.001:
                    sig_text = "*** Highly Significant"
                elif comp.p_value < 0.01:
                    sig_text = "** Very Significant"
                elif comp.p_value < 0.05:
                    sig_text = "* Significant"

                result_text = f"Column: {comp.column}  |  p-value: {comp.p_value:.4f}  |  {sig_text}"

                result_box = slide.shapes.add_textbox(
                    content_area['left'] + Inches(0.2), current_top,
                    content_area['width'] - Inches(0.4), Inches(0.35)
                )
                result_frame = result_box.text_frame
                result_frame.paragraphs[0].text = result_text
                result_frame.paragraphs[0].font.size = Pt(12)

                current_top += Inches(0.6)

        # Difference Analysis
        for diff in report_data.differences[:2]:  # 최대 2개
            title = f"차이 분석: {diff.dataset_a_name} vs {diff.dataset_b_name}" if is_ko else \
                    f"Difference: {diff.dataset_a_name} vs {diff.dataset_b_name}"
            slide, content_area = self._add_content_slide(prs, title, options)

            # 차이 바 차트 시뮬레이션
            bar_top = content_area['top'] + Inches(0.5)
            bar_height = Inches(0.6)
            bar_width = content_area['width'] - Inches(1)
            bar_left = content_area['left'] + Inches(0.5)

            # Positive bar
            pos_width = bar_width * (diff.positive_percentage / 100)
            if pos_width > 0:
                pos_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    bar_left, bar_top,
                    pos_width, bar_height
                )
                pos_bar.fill.solid()
                pos_bar.fill.fore_color.rgb = RGBColor(34, 197, 94)
                pos_bar.line.fill.background()

            # Negative bar
            neg_width = bar_width * (diff.negative_percentage / 100)
            if neg_width > 0:
                neg_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    bar_left + pos_width, bar_top,
                    neg_width, bar_height
                )
                neg_bar.fill.solid()
                neg_bar.fill.fore_color.rgb = RGBColor(239, 68, 68)
                neg_bar.line.fill.background()

            # Neutral bar
            neu_width = bar_width * (diff.neutral_percentage / 100)
            if neu_width > 0:
                neu_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    bar_left + pos_width + neg_width, bar_top,
                    neu_width, bar_height
                )
                neu_bar.fill.solid()
                neu_bar.fill.fore_color.rgb = RGBColor(148, 163, 184)
                neu_bar.line.fill.background()

            # 범례
            legend_top = bar_top + Inches(0.8)
            legend_text = (
                f"Positive (A > B): {diff.positive_percentage:.1f}%  |  "
                f"Negative (A < B): {diff.negative_percentage:.1f}%  |  "
                f"No Change: {diff.neutral_percentage:.1f}%"
            )

            legend_box = slide.shapes.add_textbox(
                bar_left, legend_top,
                bar_width, Inches(0.4)
            )
            legend_frame = legend_box.text_frame
            legend_frame.paragraphs[0].text = legend_text
            legend_frame.paragraphs[0].font.size = Pt(12)
            legend_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 요약 통계
            stats_top = legend_top + Inches(0.7)
            stats_text = (
                f"Total Difference: {self.format_number(diff.total_difference)}  |  "
                f"Mean Difference: {self.format_number(diff.mean_difference)}"
            )

            stats_box = slide.shapes.add_textbox(
                bar_left, stats_top,
                bar_width, Inches(0.4)
            )
            stats_frame = stats_box.text_frame
            stats_frame.paragraphs[0].text = stats_text
            stats_frame.paragraphs[0].font.size = Pt(14)
            stats_frame.paragraphs[0].font.bold = True
            stats_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_table_slides(
        self,
        prs: "Presentation",
        report_data: ReportData,
        options: ReportOptions
    ):
        """Table 슬라이드들"""
        for table_data in report_data.tables[:3]:  # 최대 3개 테이블
            slide, content_area = self._add_content_slide(prs, table_data.title, options)

            # 테이블 크기 계산
            num_cols = min(len(table_data.columns), 6)
            num_rows = min(len(table_data.rows), 8) + 1  # 헤더 포함

            table_width = content_area['width'] - Inches(0.5)
            row_height = Inches(0.4)

            table = slide.shapes.add_table(
                num_rows, num_cols,
                content_area['left'] + Inches(0.25),
                content_area['top'] + Inches(0.3),
                table_width,
                row_height * num_rows
            ).table

            # 헤더
            for i, col in enumerate(table_data.columns[:num_cols]):
                cell = table.cell(0, i)
                cell.text = col[:12]
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._hex_to_rgb(self.template.primary_color)

                para = cell.text_frame.paragraphs[0]
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.size = Pt(10)
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER

            # 데이터
            for row_idx, row_data in enumerate(table_data.rows[:8]):
                for col_idx, cell_value in enumerate(row_data[:num_cols]):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell_text = str(cell_value) if cell_value is not None else "-"
                    if len(cell_text) > 12:
                        cell_text = cell_text[:10] + ".."
                    cell.text = cell_text

                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(9)
                    para.alignment = PP_ALIGN.CENTER

    def _add_closing_slide(
        self,
        prs: "Presentation",
        report_data: ReportData,
        options: ReportOptions
    ):
        """마무리 슬라이드"""
        is_ko = options.language == 'ko'

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        width = prs.slide_width
        height = prs.slide_height

        # 배경
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self._hex_to_rgb(self.template.primary_color)
        bg_shape.line.fill.background()

        # 감사 메시지
        thank_text = "감사합니다" if is_ko else "Thank You"
        thank_box = slide.shapes.add_textbox(
            Inches(0.5),
            (height - Inches(2)) / 2,
            width - Inches(1),
            Inches(1)
        )
        thank_frame = thank_box.text_frame
        thank_frame.paragraphs[0].text = thank_text
        thank_frame.paragraphs[0].font.size = Pt(48)
        thank_frame.paragraphs[0].font.bold = True
        thank_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        thank_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 부가 정보
        info_text = f"Generated by Data Graph Studio\n{report_data.metadata.created_at.strftime('%Y-%m-%d')}"
        info_box = slide.shapes.add_textbox(
            Inches(0.5),
            height - Inches(1.5),
            width - Inches(1),
            Inches(0.8)
        )
        info_frame = info_box.text_frame
        info_frame.paragraphs[0].text = info_text
        info_frame.paragraphs[0].font.size = Pt(12)
        info_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)
        info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
